#!/usr/bin/env python3
"""make_deck.py — xr-splat decoupled 발표 .pptx 생성 (기승전결, 실측 수치).
usage: python make_deck.py  → docs/presentation/xr-splat-decoupled.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
STILL = ROOT / "docs/assets/reloc_demo_still_mcmc2m.png"
OUT = ROOT / "docs/presentation/xr-splat-decoupled.pptx"

NAVY = RGBColor(0x12, 0x1B, 0x2E); INK = RGBColor(0x1A, 0x1A, 0x1A)
ACC = RGBColor(0x2E, 0x86, 0xC1); GREEN = RGBColor(0x1E, 0x8E, 0x3E)
GREY = RGBColor(0x55, 0x55, 0x55); WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height


def tb(slide, l, t, w, h):
    box = slide.shapes.add_textbox(l, t, w, h); tf = box.text_frame; tf.word_wrap = True
    return tf


def para(tf, text, size, color=INK, bold=False, bullet=False, first=False, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = ("•  " if bullet else "") + text
    if align: p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Malgun Gothic"
    return p


def bar(slide, tag):
    b = slide.shapes.add_textbox(Inches(11.3), Inches(0.25), Inches(1.8), Inches(0.4))
    para(b.text_frame, tag, 12, ACC, bold=True, first=True, align=PP_ALIGN.RIGHT)


def content(title, tag, lines):
    s = prs.slides.add_slide(BLANK)
    para(tb(s, Inches(0.6), Inches(0.45), Inches(10.5), Inches(1.0)), title, 30, NAVY, bold=True, first=True)
    bar(s, tag)
    sl = s.shapes.add_shape;
    line = s.shapes.add_textbox(Inches(0.62), Inches(1.35), Inches(3), Inches(0.05))
    tf = tb(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(5.6))
    first = True
    for ln in lines:
        lvl, txt, *rest = ln
        color = rest[0] if rest else INK
        bold = lvl == 0
        para(tf, txt, 22 - lvl*3, color, bold=bold, bullet=(lvl > 0), first=first)
        first = False
    return s


# Slide 1 — 표지
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, W, H); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
para(tb(s, Inches(1.0), Inches(2.5), Inches(11.3), Inches(1.5)),
     "XR 실사급 공간 자산을 위한 SLAM × Gaussian Splatting", 36, WHITE, bold=True, first=True)
para(tb(s, Inches(1.0), Inches(4.0), Inches(11.3), Inches(1.2)),
     "통합형의 한계와 decoupled 구조로의 전환 — 설계부터 런타임 실증까지", 20, RGBColor(0xBB, 0xCC, 0xEE), first=True)

content("목표", "起", [
    (0, "“XR에서 사람이 공간에 들어가도 어색하지 않은 실사급 공간 자산”", NAVY),
    (1, "보이는 화면 = 실사 같은 3D 공간"),
    (1, "동시에 = 사용자 위치를 안정적으로 추정 (localization)"),
    (1, "핵심 질문: 이걸 한 모델로? 아니면 나눠서?", ACC)])

content("첫 시도 — 통합형 Gaussian-SLAM", "起", [
    (0, "하나의 Gaussian 표현으로 tracking · mapping · rendering을 실시간 동시 수행", NAVY),
    (1, "실험: SplaTAM · GS-SLAM · Photo-SLAM · LoopSplat 4종 직접 빌드·실행"),
    (1, "기대: 한 시스템으로 위치추정 + 실사 렌더 다 해결")])

content("통합형의 한계", "承", [
    (0, "4종 실험 결과 — 세 가지가 동시에 미달", NAVY),
    (1, "렌더 품질: 실사급 미달", RGBColor(0xC0,0x39,0x2B)),
    (1, "트래킹: 성숙한 전용 SLAM보다 약함", RGBColor(0xC0,0x39,0x2B)),
    (1, "속도: 무거움 (실시간 XR 부담)", RGBColor(0xC0,0x39,0x2B)),
    (1, "→ ‘장점 통합’이 아니라 ‘단점 통합’ 위험", ACC)])

content("전환 결정 — decoupled (분리)", "轉", [
    (0, "역할을 쪼갠다", NAVY),
    (1, "Localization → ORB-SLAM3 (성숙한 포즈·루프클로저·relocalization)"),
    (1, "Rendering → gsplat (오프라인 무제약 실사 최적화)"),
    (1, "핵심: Gaussian 맵을 SLAM 포즈 위에서 학습 → 동일 좌표계 → 정합 0단계", GREEN)])

content("decoupled 파이프라인", "轉", [
    (0, "오프라인 자산 생성 + 런타임 위치추정 분리", NAVY),
    (1, "[오프라인] RGB-D → ORB-SLAM3 포즈 → COLMAP → gsplat 학습(포즈고정+depth) → 경량 자산"),
    (1, "[런타임] 새 프레임 → 저장된 맵에 relocalize → 그 pose로 Gaussian 렌더"),
    (1, "8단계 독립 CLI, 표준 포맷(TUM/COLMAP)으로 통신 → 단계 교체 자유")])

content("M1 — 원리 검증 (공개 데이터)", "轉", [
    (0, "“SLAM 포즈로도 실사급 GS가 나오나?” — TUM fr1/desk, hold-out 16뷰", NAVY),
    (1, "ORB 포즈 23.85  vs  COLMAP 포즈 23.98 PSNR  →  차이 0.13dB", GREEN),
    (1, "ATE 둘 다 ~2cm"),
    (1, "✅ decoupled 원리 성립 — 좋은 입력이면 SLAM 포즈로 COLMAP급 품질", ACC)])

content("M2 — 실제 D455 캡처", "轉", [
    (0, "자체 캡처 3개 — 캡처 기하가 품질을 좌우", NAVY),
    (1, "room1: 궤적 0.43m (제자리회전, narrow)"),
    (1, "room2: 궤적 0.65m (좁은 영역)"),
    (1, "home: 궤적 28.78m (공간 횡단, baseline 충분)", GREEN),
    (1, "교훈: 걸으면서 시차(parallax) 확보가 실사 품질의 전제", ACC)])

content("난관 1 — 포즈 품질이 병목", "轉", [
    (0, "room1에서 발견 — 같은 캡처인데 포즈를 바꾸니 +5dB", NAVY),
    (1, "ORB 포즈 vs full SfM 포즈: per-pose 10cm 차이"),
    (1, "SfM 포즈로 학습 시 +5dB 선명", GREEN),
    (1, "→ 실제 캡처에선 ORB 포즈 정확도가 실사 품질의 병목", ACC)])

content("난관 2 — 프레임 정합 함정 (핵심 발견)", "轉", [
    (0, "“PSNR 높다고 좋은 자산이 아니다”", NAVY),
    (1, "home은 PSNR 24.34로 최고였지만 ORB 프레임에서 76° 틀어짐 → 런타임 위치추정 불가", RGBColor(0xC0,0x39,0x2B)),
    (1, "해결: snap_scene_to_orb.py로 ORB 프레임에 강체 재정렬 → scale 1.0 / rot 0°"),
    (1, "결과: home 23.82 PSNR + 프레임 유효 (−0.52dB로 좌표계 회복) — 재캡처 0", GREEN)])

# Slide 11 — 런타임 실증 (이미지 임베드)
s = prs.slides.add_slide(BLANK)
para(tb(s, Inches(0.6), Inches(0.45), Inches(11), Inches(1)), "런타임 reloc → render 실증", 30, NAVY, bold=True, first=True)
bar(s, "轉→結")
tf = tb(s, Inches(0.7), Inches(1.5), Inches(6.0), Inches(5.4))
first = True
for lvl, txt, *c in [
    (0, "decoupled의 존재 이유를 end-to-end로 증명", NAVY),
    (1, "저장된 맵에 새 프레임 relocalize (COLMAP PnP)"),
    (1, "그 pose로 Gaussian 렌더"),
    (1, "reloc 성공 10 / 10 (100%)", GREEN),
    (1, "좌표 앵커 오차 0.000cm", GREEN),
    (1, "실제 화면 = 가우시안 렌더 (TV·식물·소파·창 동일 위치)", ACC)]:
    para(tf, txt, 20 - lvl*2, (c[0] if c else INK), bold=(lvl == 0), bullet=(lvl > 0), first=first); first = False
if STILL.exists():
    s.shapes.add_picture(str(STILL), Inches(6.9), Inches(2.3), width=Inches(6.0))
    para(tb(s, Inches(6.9), Inches(4.6), Inches(6.0), Inches(0.4)), "좌: 실제  |  우: 가우시안(reloc pose)", 13, GREY, first=True, align=PP_ALIGN.CENTER)

content("결과 종합", "結", [
    (0, "통합형 → decoupled, 처음부터 끝까지 데이터로 증명", NAVY),
    (1, "통합형 4종: 한계 확인 (렌더·트래킹·속도)"),
    (1, "decoupled 원리(M1): ✅ Δ0.13dB", GREEN),
    (1, "실사 자산(M2 home, mcmc2m): ✅ PSNR 27.96 / SSIM 0.871 / LPIPS 0.267, 프레임 유효", GREEN),
    (1, "런타임 reloc→render: ✅ 10/10, 앵커 0cm", GREEN),
    (1, "→ ‘통합 모델보다 분리가 옳았다’", ACC)])

content("품질 해소 & 다음", "結", [
    (0, "시각 품질 — MCMC strategy로 해소 (mcmc2m 채택)", NAVY),
    (1, "기존 2M(refine-stop 5k 동결) 23.82 → mcmc2m 27.96 PSNR (+4.1dB, LPIPS 0.39→0.27)", GREEN),
    (1, "병목은 raw-count가 아니라 배치/정제 — 같은 2M을 계속 정제하니 풀림"),
    (1, "scene 타일링: 통제 비교(2M≈3M)로 불필요 확인 → 선반행", ACC),
    (0, "다음: ad-server 복귀 (motion 18개 smoke test)", NAVY)])

content("결론", "結", [
    (0, "통합형 Gaussian-SLAM은 이론은 매력적이나 렌더·트래킹·속도를 동시 만족 못 함", INK),
    (0, "SLAM(위치) + Gaussian(렌더)을 분리하고 동일 좌표계로 연동하면,", NAVY),
    (0, "안정적 localization 위에서 실사 공간을 렌더하는 구조가 실제로 작동한다", NAVY),
    (1, "설계 → 자산 → 런타임까지 전 구간 실증 완료", GREEN)])

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"[deck] {len(prs.slides.__iter__.__self__._sldIdLst)} slides → {OUT}")
