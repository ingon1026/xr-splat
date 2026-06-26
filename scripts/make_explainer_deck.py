#!/usr/bin/env python3
"""make_explainer_deck.py — xr-splat 전체 흐름 A to Z 설명 덱 (다이어그램·수식·실측 이미지).
usage: python scripts/make_explainer_deck.py → docs/presentation/xr-splat-explainer.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "docs/assets/report"
OUT = ROOT / "docs/presentation/xr-splat-explainer.pptx"

NAVY = RGBColor(0x12, 0x1B, 0x2E); INK = RGBColor(0x20, 0x20, 0x20)
BLUE = RGBColor(0x2E, 0x86, 0xC1); GREEN = RGBColor(0x1E, 0x8E, 0x3E)
ORANGE = RGBColor(0xE6, 0x7E, 0x22); PURPLE = RGBColor(0x7D, 0x3C, 0x98)
GREY = RGBColor(0x88, 0x88, 0x88); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGREY = RGBColor(0xEC, 0xF0, 0xF1)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; W, H = prs.slide_width, prs.slide_height


def slide():
    return prs.slides.add_slide(BLANK)


def tb(s, l, t, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.text = line; p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Malgun Gothic"
    return box


def title(s, text, tag=""):
    tb(s, 0.55, 0.35, 11.5, 0.9, text, 30, NAVY, bold=True)
    if tag:
        tb(s, 11.0, 0.42, 1.9, 0.5, tag, 13, BLUE, bold=True, align=PP_ALIGN.RIGHT)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25), Inches(12.1), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()


def box(s, l, t, w, h, text, fill, fg=WHITE, size=14, bold=True):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = WHITE; sh.line.width = Pt(1)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.text = line; p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(size - (2 if i else 0)); r.font.bold = bold if i == 0 else False
            r.font.color.rgb = fg; r.font.name = "Malgun Gothic"
    return sh


def arrow(s, l, t, w, h=0.5, color=GREY, label=""):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()
    if label:
        tb(s, l - 0.1, t - 0.45, w + 0.2, 0.4, label, 11, color, bold=True, align=PP_ALIGN.CENTER)


def pic(s, path, l, t, w, cap=""):
    if Path(path).exists():
        p = s.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))
        if cap:
            tb(s, l, t + Emu(p.height).inches + 0.05, w, 0.4, cap, 11, GREY, align=PP_ALIGN.CENTER)
        return p


# ---------- Slide 1: 표지 ----------
s = slide()
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tb(s, 1.0, 2.4, 11.3, 1.4, "xr-splat — A to Z 전체 흐름", 40, WHITE, bold=True)
tb(s, 1.0, 3.9, 11.3, 1.0, "카메라 입력부터 화면 출력까지, 그림으로 쉽게", 20, RGBColor(0xBB, 0xCC, 0xEE))
tb(s, 1.0, 5.1, 11.3, 0.6, "SLAM(위치) × Gaussian Splatting(렌더) decoupled 파이프라인", 15, RGBColor(0x88, 0xAA, 0xCC))

# ---------- Slide 2: 한 장 답 (토스토스 vs 합치기) ----------
s = slide(); title(s, "먼저 핵심: 따로 만들어 합치는 게 아니라 '순차(토스토스)'", "답")
tb(s, 0.6, 1.5, 12.1, 0.7, "SLAM과 Gaussian을 따로 만들어 나중에 억지로 붙이는 게 아님. 순서대로 흘리되,", 17, INK)
tb(s, 0.6, 2.05, 12.1, 0.7, "Gaussian을 'SLAM이 만든 카메라 포즈 위에서' 학습 → 자동으로 같은 좌표계가 됨.", 17, INK, bold=True)
box(s, 0.8, 3.1, 2.4, 1.2, "① D455\nRGB + Depth", BLUE)
arrow(s, 3.35, 3.5, 0.8, 0.45, GREY)
box(s, 4.25, 3.1, 2.4, 1.2, "② SLAM\n카메라 위치", GREEN)
arrow(s, 6.8, 3.5, 0.8, 0.45, GREY)
box(s, 7.7, 3.1, 2.4, 1.2, "③ Gaussian 학습\n(②의 포즈 위에서)", ORANGE)
arrow(s, 10.25, 3.5, 0.8, 0.45, GREY)
box(s, 11.15, 3.1, 1.9, 1.2, "④ 실사 자산\n.ply", PURPLE)
tb(s, 0.8, 4.7, 12.2, 1.6,
   "→ ③에서 Gaussian이 ②의 좌표계를 그대로 물려받음. 그래서 나중에 '어디 있나(위치추정)'를\n"
   "   풀면 그 좌표를 변환 없이 Gaussian 렌더러에 바로 넣을 수 있음. 이게 decoupled의 핵심 트릭.",
   16, GREEN, bold=True)

# ---------- Slide 3: 입력 D455 ----------
s = slide(); title(s, "① 입력 — D455 카메라가 주는 것", "입력")
box(s, 1.0, 1.7, 3.6, 1.3, "RGB 영상\n(색)", BLUE, size=16)
box(s, 1.0, 3.2, 3.6, 1.3, "Depth 영상\n(픽셀마다 거리, m)", GREEN, size=16)
tb(s, 5.0, 1.7, 7.8, 3.5,
   "• D455 = RGB-D 카메라. 한 프레임마다 두 장: 색(RGB) + 깊이(Depth).\n\n"
   "• Depth가 핵심: 픽셀이 '얼마나 멀리'인지 알면 2D 점을 3D로 띄울 수 있음.\n"
   "  (특징점을 3D로 backproject → 나중에 위치추정·맵에 사용)\n\n"
   "• 실측 D455 캡처를 그대로 사용 (공개 데이터 아님). 예: home, room2.\n\n"
   "• 한 번 녹화(.bag) → 오프라인으로 추출 (카메라 다시 안 켜도 됨).",
   16, INK)

# ---------- Slide 4: 오프라인 파이프라인 (토스토스 단계별) ----------
s = slide(); title(s, "오프라인: 자산 만들기 (토스토스 4단계)", "파이프라인")
ys = 1.7
steps = [
    ("①\nD455 .bag", "RGB+Depth\n프레임 추출", BLUE),
    ("②\nORB-SLAM3", "카메라가 어디\n있었나 = 포즈", GREEN),
    ("③\nCOLMAP", "포즈 정밀화\n+ 3D 점", RGBColor(0x16,0xA0,0x85)),
    ("④\ngsplat 학습", "③ 포즈 고정하고\nGaussian 최적화", ORANGE),
    ("⑤\n자산 .ply", "실사 Gaussian\n맵 완성", PURPLE),
]
x = 0.55
for i, (t, sub, c) in enumerate(steps):
    box(s, x, ys, 2.05, 1.5, t + "\n" + sub, c, size=14)
    x += 2.05
    if i < len(steps) - 1:
        arrow(s, x, ys + 0.5, 0.35, 0.45, GREY); x += 0.35
tb(s, 0.6, 3.6, 12.2, 2.6,
   "각 화살표 = '앞 단계 결과를 파일로 다음 단계에 토스'. 표준 포맷(TUM/COLMAP)으로 통신.\n\n"
   "• ② SLAM: 영상만 보고 카메라 궤적·키프레임을 추정 (위치의 기준 좌표계 생성).\n"
   "• ③ COLMAP: 그 포즈를 다듬고, depth로 3D 점을 만들어 Gaussian 초기값 제공.\n"
   "• ④ 학습: 포즈는 '고정'하고 Gaussian만 움직임 → Gaussian이 SLAM 좌표계에 박힘.",
   15, INK)

# ---------- Slide 5: 왜 합쳐지나 (같은 좌표계) ----------
s = slide(); title(s, "왜 '합쳐진 맵'이 되나 — 같은 좌표계", "핵심")
box(s, 1.2, 1.8, 4.2, 1.2, "SLAM 좌표계\n(카메라가 어디 있나)", GREEN, size=16)
a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(3.0), Inches(3.1), Inches(0.6), Inches(0.9))
a.fill.solid(); a.fill.fore_color.rgb = GREY; a.line.fill.background()
tb(s, 3.8, 3.25, 6.0, 0.6, "이 포즈를 '고정'하고 그 위에서 학습", 14, GREY, bold=True)
box(s, 1.2, 4.1, 4.2, 1.2, "Gaussian 맵\n(같은 좌표계에 박힘)", ORANGE, size=16)
tb(s, 6.0, 1.8, 6.8, 4.4,
   "보통은 SLAM 맵과 렌더 맵을 따로 만들면\n나중에 정합(align)이 지옥.\n\n"
   "여기선 Gaussian을 '처음부터 SLAM 포즈 위에서'\n학습 → 두 맵이 같은 좌표계.\n\n"
   "그래서 런타임에 위치추정으로 얻은 포즈(Tcw)를\n변환 0단계로 Gaussian 렌더러에 바로 투입.\n\n"
   "= 이게 'decoupled SLAM × Gaussian'의 존재 이유.",
   16, INK)

# ---------- Slide 6: 학습 수식 (쉽게) ----------
s = slide(); title(s, "④ 학습은 무엇을 맞추나 (수식, 쉽게)", "학습")
tb(s, 0.6, 1.5, 12.1, 0.7, "Gaussian 맵 = 공간에 떠 있는 수백만 개의 '색칠된 타원 구름'. 각 Gaussian의 파라미터:", 16, INK)
box(s, 0.7, 2.3, 12.0, 0.9, "위치 μ   |   크기·방향 Σ   |   투명도 α   |   색(SH 계수)", LGREY, INK, size=16, bold=True)
tb(s, 0.6, 3.5, 12.1, 0.6, "학습: 이 Gaussian들을 렌더한 그림이 실제 사진과 같아지도록 파라미터를 경사하강으로 조정.", 16, INK)
box(s, 1.5, 4.4, 10.3, 1.0,
    "Loss = (1−λ)·|렌더−실제|  +  λ·(1−SSIM)  +  λd·|렌더깊이−센서깊이|", NAVY, WHITE, size=18, bold=True)
tb(s, 0.6, 5.7, 12.1, 1.2,
   "• 1항: 색이 실제와 같게   • 2항(SSIM): 구조·또렷함   • 3항: 깊이도 센서와 맞게(형태 안정)\n"
   "• 포즈는 고정 → Gaussian만 학습 (그래서 SLAM 좌표계 유지). MCMC 전략으로 배치 최적화.",
   15, GREEN)

# ---------- Slide 7: 런타임 ----------
s = slide(); title(s, "런타임: 새로 들어왔을 때 (위치추정 → 렌더)", "런타임")
box(s, 0.7, 2.0, 2.7, 1.3, "새 카메라\n프레임 (RGB)", BLUE, size=15)
arrow(s, 3.55, 2.4, 0.7, 0.45, GREEN, "위치 찾기")
box(s, 4.45, 2.0, 3.0, 1.3, "PnP / photometric\n위치추정\n→ 포즈 Tcw", GREEN, size=14)
arrow(s, 7.65, 2.4, 0.7, 0.45, ORANGE, "변환 0")
box(s, 8.55, 2.0, 2.6, 1.3, "3DGS 렌더\n(그 포즈로)", ORANGE, size=15)
arrow(s, 11.3, 2.4, 0.6, 0.45, PURPLE)
box(s, 11.2, 3.5, 1.9, 1.0, "화면\n실사", PURPLE, size=15)
a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(11.95), Inches(3.35), Inches(0.4), Inches(0.2))
a.fill.solid(); a.fill.fore_color.rgb = PURPLE; a.line.fill.background()
tb(s, 0.6, 3.7, 10.3, 3.0,
   "• 위치추정 = '지금 카메라가 맵의 어디인가'를 이미지만 보고 푸는 것.\n"
   "   - feature-PnP: 특징점 매칭 → solvePnP, ~27 FPS, 전역(아무데서나) 가능.\n"
   "   - photometric: Gaussian 렌더를 이미지에 맞춰 포즈 미세조정.\n\n"
   "• 얻은 포즈는 Gaussian과 같은 좌표계 → 변환 없이 렌더 → 화면.\n\n"
   "• 즉 '실제 카메라 → 위치 → 그 자리 실사' 가 한 루프.",
   15, INK)

# ---------- Slide 8: 아웃풋이 어떻게 보이나 (localize→render) ----------
s = slide(); title(s, "출력 1 — 찾은 포즈로 그린 게 실제와 일치", "출력")
pic(s, IMG / "localize_to_render.png", 0.6, 1.5, 6.0, "왼쪽=실제 query / 오른쪽=PnP가 찾은 포즈로 렌더 (27~29 dB)")
tb(s, 7.0, 1.7, 6.0, 4.8,
   "각 행:\n실제 사진(왼쪽) → 시스템이 그 사진만 보고\n위치를 찾음 → 그 포즈로 Gaussian 렌더(오른쪽).\n\n"
   "오른쪽은 '정답 포즈'를 준 게 아니라\nPnP가 찾아낸 포즈로만 그린 것.\n\n"
   "그런데 실제와 27~29 dB로 일치\n→ 위치추정 + 렌더가 한 좌표계에서\n   정확히 맞물린다는 end-to-end 증거.",
   16, INK)

# ---------- Slide 9: SLAM 경로 위 위치 ----------
s = slide(); title(s, "출력 2 — 위치가 SLAM 경로 위 올바른 자리에", "출력")
pic(s, IMG / "localization_on_slam_path.png", 0.6, 1.5, 7.2)
tb(s, 8.1, 1.8, 4.9, 4.6,
   "초록 = SLAM이 추정한 카메라 경로.\n빨강 = 위치추정이 찾은 query 카메라\n(화살표 = 보는 방향).\n\n"
   "빨간 점이 초록 경로 위에 정확히 얹힘\n→ 렌더뿐 아니라 '위치'도 맞다.\n\n"
   "공간적으로 시스템이 제대로 작동.",
   16, INK)

# ---------- Slide 10: 커버리지 (네 질문: 맵 전체로 안 되나) ----------
s = slide(); title(s, "'맵 전체로 되잖아?' — 맞아, 학습한 만큼 커버", "커버리지")
tb(s, 0.6, 1.5, 12.2, 0.8, "지금 home 자산은 전체 29m 녹화 중 '한 구간(2m 영역)'만 학습한 것. 한계가 아니라 처리 선택.", 17, INK, bold=True)
box(s, 0.9, 2.6, 5.6, 1.0, "전체 녹화: 29m 경로, 8.7m 공간", GREY, WHITE, size=15)
box(s, 0.9, 3.8, 5.6, 1.0, "현재 학습 자산: 2m 영역 (그 일부)", ORANGE, WHITE, size=15)
tb(s, 6.9, 2.6, 6.0, 4.0,
   "• Gaussian은 '학습에 넣은 프레임'만큼만 커버.\n"
   "• 전체 29m를 다 학습하면 전체 공간 자산이 됨\n  (COLMAP+학습을 전 구간에 돌리면 끝 = 처리량 문제).\n\n"
   "• 그래서 '어디든 걸어가기'엔 더 넓게 학습/캡처하면 됨.\n  품질·범위는 '얼마나 찍고 학습했나'가 결정.",
   16, INK)

# ---------- Slide 11: 품질은 캡처가 좌우 (home vs room2) ----------
s = slide(); title(s, "품질은 '어떻게 찍었나'가 좌우", "비교")
pic(s, IMG / "home_vs_room2.png", 0.6, 1.6, 7.4, "위=home(넓게 둘러봄, 또렷) / 아래=room2(좁게, soft)")
tb(s, 8.3, 1.8, 4.7, 4.4,
   "같은 시스템, 다른 캡처:\n\n"
   "• home: 시차(parallax) 충분 → 선명.\n• room2: 좁은 캡처 → soft.\n\n"
   "→ 알고리즘이 아니라 '캡처 시차/범위'가\n   실사 품질의 천장.",
   16, INK)

# ---------- Slide 12: 한 장 요약 ----------
s = slide(); title(s, "한 장 요약 — A to Z", "요약")
flow = [
    ("D455\nRGB+Depth", BLUE), ("SLAM\n포즈", GREEN), ("COLMAP\n포즈+점", RGBColor(0x16,0xA0,0x85)),
    ("Gaussian 학습\n(포즈 고정)", ORANGE), ("실사 자산", PURPLE),
]
x = 0.7
for i, (t, c) in enumerate(flow):
    box(s, x, 1.8, 2.1, 1.1, t, c, size=13); x += 2.1
    if i < len(flow) - 1:
        arrow(s, x, 2.15, 0.3, 0.4, GREY); x += 0.3
tb(s, 0.7, 3.2, 12.2, 0.5, "── 런타임 ──", 14, GREY, bold=True, align=PP_ALIGN.CENTER)
x = 1.6
run = [("새 프레임", BLUE), ("위치추정\n(PnP, 27FPS)", GREEN), ("같은 좌표계\n변환 0", GREY), ("3DGS 렌더", ORANGE), ("화면 실사", PURPLE)]
for i, (t, c) in enumerate(run):
    box(s, x, 3.7, 2.0, 1.0, t, c, size=13); x += 2.0
    if i < len(run) - 1:
        arrow(s, x, 4.0, 0.3, 0.4, GREY); x += 0.3
tb(s, 0.7, 5.1, 12.2, 1.6,
   "핵심 한 줄: D455로 찍어 → SLAM이 좌표계를 잡고 → 그 위에서 Gaussian을 학습해 실사 맵을 만들고\n"
   "→ 런타임엔 새 프레임의 위치를 찾아(같은 좌표계) → 그 자리에서 Gaussian을 렌더해 실사로 보여준다.",
   16, GREEN, bold=True)

OUT.parent.mkdir(parents=True, exist_ok=True); prs.save(str(OUT))
print(f"[explainer] {len(prs.slides._sldIdLst)} slides → {OUT}")
