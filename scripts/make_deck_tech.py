#!/usr/bin/env python3
"""make_deck_tech.py — xr-splat 기술 원리 설명 .pptx (알고리즘·맵생성·렌더링).
대상: 비전 엔지니어(ML/CV는 알지만 SLAM/3DGS는 처음). usage: python make_deck_tech.py"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs/presentation/xr-splat-technical.pptx"
NAVY = RGBColor(0x12, 0x1B, 0x2E); INK = RGBColor(0x1A, 0x1A, 0x1A)
ACC = RGBColor(0x2E, 0x86, 0xC1); GREEN = RGBColor(0x1E, 0x8E, 0x3E)
ORANGE = RGBColor(0xB9, 0x6A, 0x00); GREY = RGBColor(0x55, 0x55, 0x55); WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; W, H = prs.slide_width, prs.slide_height


def tb(s, l, t, w, h):
    box = s.shapes.add_textbox(l, t, w, h); box.text_frame.word_wrap = True; return box.text_frame


def para(tf, text, size, color=INK, bold=False, lvl=0, first=False, align=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = ("   " * (lvl - 1)) + ("•  " if lvl else "") + text
    if align: p.alignment = align
    for r in p.runs:
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Malgun Gothic"
    return p


def content(title, tag, lines):
    s = prs.slides.add_slide(BLANK)
    para(tb(s, Inches(0.6), Inches(0.42), Inches(11), Inches(0.9)), title, 28, NAVY, bold=True, first=True)
    para(tb(s, Inches(11.3), Inches(0.25), Inches(1.8), Inches(0.4)), tag, 12, ACC, bold=True, first=True, align=PP_ALIGN.RIGHT)
    tf = tb(s, Inches(0.7), Inches(1.45), Inches(12.0), Inches(5.7)); first = True
    for ln in lines:
        lvl, txt, *c = ln
        para(tf, txt, 21 - min(lvl, 2) * 3, (c[0] if c else INK), bold=(lvl == 0), lvl=lvl, first=first); first = False
    return s


# 1 표지
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(1, 0, 0, W, H); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
para(tb(s, Inches(1), Inches(2.4), Inches(11.3), Inches(1.6)), "맵은 어떻게 만들고, 렌더링은 어떻게 되는가", 36, WHITE, bold=True, first=True)
para(tb(s, Inches(1), Inches(4.0), Inches(11.3), Inches(1)), "xr-splat 기술 원리 — ORB-SLAM3 · 3D Gaussian Splatting · PnP relocalization", 20, RGBColor(0xBB, 0xCC, 0xEE), first=True)

# 2 전체 그림
content("전체 그림 — 3개 알고리즘이 각자 일한다", "개요", [
    (0, "맵 생성(오프라인)", ACC),
    (1, "① ORB-SLAM3 → 카메라 포즈(위치) 추정 + 희소 맵"),
    (1, "② COLMAP SfM → 포즈 정밀화 + 3D 점(매칭 인프라)"),
    (1, "③ 3D Gaussian Splatting → 그 포즈 위에서 실사 렌더용 맵 학습"),
    (0, "런타임", ACC),
    (1, "④ PnP relocalization → 새 프레임의 위치를 맵에서 찾음 → 그 위치로 ③ 렌더"),
    (0, "핵심: ①②가 '어디(geometry)', ③이 '어떻게 보이는지(appearance)' 담당", GREEN)])

# 3 ORB-SLAM3 개요
content("① ORB-SLAM3 — 위치추정(SLAM)", "알고리즘", [
    (0, "SLAM = 카메라가 움직이며 동시에 ‘내 위치’ + ‘주변 지도’를 추정", NAVY),
    (1, "입력: RGB-D 영상 시퀀스 (depth가 있어 metric 스케일 자동 확보)"),
    (1, "특징점 기반(feature-based): 픽셀이 아니라 ORB 특징점으로 기하 계산"),
    (1, "3개 스레드가 병렬로 동작 (다음 장)"),
    (1, "출력: 키프레임별 카메라 포즈 Tcw + 희소 3D 맵 → Atlas(.osa)로 저장", GREEN)])

# 4 ORB feature + 3 threads
content("① ORB-SLAM3 — ORB 특징점과 3 스레드", "알고리즘", [
    (0, "ORB = Oriented FAST + Rotated BRIEF", ACC),
    (1, "FAST: 코너 검출(빠름) / BRIEF: 이진 디스크립터 → Hamming 거리로 초고속 매칭"),
    (1, "회전·스케일 불변 → 시점 바뀌어도 같은 점 매칭"),
    (0, "3 스레드", ACC),
    (2, "Tracking: 프레임마다 이전 맵에 특징 매칭 → 모션모델+PnP로 포즈 추정"),
    (2, "Local Mapping: 키프레임 추가, 새 점 삼각측량, local Bundle Adjustment"),
    (2, "Loop Closing: DBoW2(bag-of-words) 장소 재인식 → 루프 검출 → pose graph 최적화")])

# 5 COLMAP
content("② COLMAP SfM — 포즈 정밀화 & 매칭 인프라", "알고리즘", [
    (0, "SfM(Structure-from-Motion) = 여러 사진에서 카메라 포즈 + 3D 점 복원", NAVY),
    (1, "SIFT 특징 추출 → 이미지 간 매칭(exhaustive/sequential)"),
    (1, "Bundle Adjustment: 재투영 오차 최소화로 포즈·점·내부파라미터 동시 최적화"),
    (0, "이 프로젝트에서 쓴 두 기능", ACC),
    (2, "point_triangulator: 포즈를 ‘고정’하고 점만 삼각측량 → 디스크립터 가진 맵 생성"),
    (2, "image_registrator: 새 이미지를 기존 맵에 2D-3D 매칭+PnP로 ‘등록’(=relocalization)")])

# 6 3DGS 표현
content("③ 3D Gaussian Splatting — 공간의 ‘표현’", "알고리즘", [
    (0, "공간 = 수백만 개의 3D 가우시안(타원체)들의 집합", NAVY),
    (0, "가우시안 1개가 가진 파라미터", ACC),
    (2, "μ : 3D 위치 (중심)"),
    (2, "Σ = R·S·Sᵀ·Rᵀ : 모양/크기 (회전 R[쿼터니언] + 스케일 S)"),
    (2, "α : 불투명도(opacity)"),
    (2, "SH 계수 : 색 — 보는 방향에 따라 달라짐(view-dependent)"),
    (0, "왜 가우시안? → 미분 가능 + 빠른 래스터라이즈 + 빈 공간 효율적", GREEN)])

# 7 렌더링
content("③ 렌더링 — ‘splatting’ 래스터라이즈", "렌더링", [
    (0, "3D 가우시안 → 화면에 그리는 과정 (NeRF처럼 광선적분 X, 투영 O → 빠름)", NAVY),
    (1, "1) 투영: 각 3D 가우시안을 카메라 평면에 2D 가우시안으로 사영(EWA splatting)"),
    (1, "2) 정렬: 깊이순으로 타일 단위 정렬"),
    (1, "3) α-blending: 앞→뒤로 색을 누적  C = Σ cᵢ·αᵢ·∏(1−αⱼ)"),
    (1, "4) 색 cᵢ = SH(보는 방향) → 반사·하이라이트 같은 시점 의존 효과"),
    (0, "전 과정 미분 가능 → 학습에 그대로 역전파 + 학습 후 실시간 렌더", GREEN)])

# 8 학습
content("③ 맵 학습 — 사진으로 가우시안을 ‘맞춰간다’", "맵 생성", [
    (0, "최적화 = 렌더 결과를 실제 사진에 맞도록 가우시안 파라미터를 경사하강", NAVY),
    (1, "초기화: 희소 점군(SfM/SLAM 점)에서 가우시안 시작"),
    (1, "반복: [학습뷰 렌더] → [실제 사진과 Loss] → [역전파 → Adam 업데이트]"),
    (2, "Loss = (1−λ)·L1 + λ·(1−SSIM)  +  depth L1 (RGB-D depth 감독)"),
    (1, "Adaptive Density: 부족한 곳 가우시안 분열/복제(densify), 투명한 것 제거(prune)"),
    (1, "이 프로젝트: 카메라 포즈 ‘고정’(SLAM 포즈 신뢰) + depth 감독 + 30k step", GREEN)])

# 9 좌표계
content("좌표계 — 두 맵을 같은 프레임에 두는 트릭", "연결", [
    (0, "포즈 표기: Tcw = world→camera (COLMAP/ORB 저장 형식)", NAVY),
    (1, "카메라 중심 C = −R_cwᵀ·t , 렌더러엔 view matrix(Tcw) 그대로 투입"),
    (0, "핵심 아이디어", ACC),
    (1, "Gaussian 맵을 ‘SLAM 포즈 위에서’ 학습 → SLAM 맵과 동일 좌표계"),
    (1, "→ 런타임에 SLAM이 준 포즈를 변환 없이 렌더러에 꽂으면 정합됨", GREEN),
    (1, "(이 프로젝트의 난관: 좌표계가 틀어지면 PSNR 좋아도 런타임 무용)", ORANGE)])

# 10 PnP reloc
content("④ PnP relocalization — 런타임 위치추정", "런타임", [
    (0, "새 프레임이 들어오면 ‘맵에서 내 위치’를 어떻게 찾나", NAVY),
    (1, "1) 새 프레임 특징 추출 → 맵의 3D 점들과 2D-3D 매칭"),
    (1, "2) PnP(Perspective-n-Point): n개의 2D↔3D 대응에서 카메라 포즈(R,t) 해석적으로 계산"),
    (1, "3) RANSAC으로 오매칭 제거 (보통 ≥15 대응 필요)"),
    (0, "결과 포즈가 ③ 맵과 같은 좌표계 → 그 포즈로 가우시안 렌더 → 실제와 정합", GREEN)])

# 11 맵빌드 A-Z
content("맵 생성 — A to Z", "파이프라인", [
    (0, "오프라인 (한 번)", ACC),
    (1, "A. RGB-D 캡처(.bag) → 프레임/깊이 추출"),
    (1, "B. ORB-SLAM3 → 키프레임 포즈(Tcw) + Atlas 맵"),
    (1, "C. COLMAP 변환/정밀화 → 포즈를 metric·정확하게, 디스크립터 맵"),
    (1, "D. 초기 점군 → 3D Gaussian Splatting 학습(포즈 고정 + depth)"),
    (1, "E. 후처리(가지치기·SH 축소) → 경량 .ply 자산"),
    (0, "산출: 가우시안 맵(.ply) + relocalization용 맵 — 동일 좌표계", GREEN)])

# 12 런타임 A-Z
content("런타임 렌더 — A to Z", "파이프라인", [
    (0, "사용 시 (매 프레임)", ACC),
    (1, "A. 새 카메라 프레임 입력"),
    (1, "B. 특징 추출 → 맵의 3D 점과 매칭"),
    (1, "C. PnP + RANSAC → 카메라 포즈(맵 좌표계)"),
    (1, "D. 그 포즈로 가우시안 맵 splatting 렌더"),
    (1, "E. 화면 출력 (= 실사 공간을 그 시점에서 본 모습)"),
    (0, "위치추정(B·C)과 렌더(D)가 분리 → 각자 최적 알고리즘 사용", GREEN)])

# 13 한 장 요약
content("한 장 요약 — 어떤 알고리즘이 어디서 뭘 하나", "정리", [
    (0, "ORB-SLAM3 → ‘어디(where)’: 카메라 포즈·맵 (ORB 특징, BA, DBoW2)", ACC),
    (0, "COLMAP SfM → 포즈 정밀화 + reloc 인프라 (SIFT, triangulation, PnP)", ACC),
    (0, "3D Gaussian Splatting → ‘어떻게 보이는지(appearance)’: 실사 렌더 (splatting, SH)", ACC),
    (0, "PnP relocalization → 런타임 위치추정 (2D-3D 매칭 + PnP/RANSAC)", ACC),
    (1, ""),
    (0, "맵 생성 = 기하(SLAM/SfM) + 외형(3DGS),  렌더 = 포즈 + splatting", GREEN),
    (0, "분리했기에 위치는 정확하게, 화면은 실사같게 — 각자 잘하는 걸로", NAVY)])

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"[tech-deck] {len(prs.slides._sldIdLst)} slides → {OUT}")
