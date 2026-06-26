#!/usr/bin/env python3
"""make_final_deck.py — xr-splat 최종 기승전결 덱 (사진 다수 + 실측 지표).
usage: python scripts/make_final_deck.py → docs/presentation/xr-splat-final.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parents[1]
IMG = ROOT / "docs/assets/report"
OUT = ROOT / "docs/presentation/xr-splat-final.pptx"

NAVY = RGBColor(0x12, 0x1B, 0x2E); INK = RGBColor(0x22, 0x22, 0x22)
BLUE = RGBColor(0x2E, 0x86, 0xC1); GREEN = RGBColor(0x1E, 0x8E, 0x3E)
ORANGE = RGBColor(0xE6, 0x7E, 0x22); PURPLE = RGBColor(0x7D, 0x3C, 0x98)
RED = RGBColor(0xC0, 0x39, 0x2B); GREY = RGBColor(0x88, 0x88, 0x88); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGREY = RGBColor(0xEC, 0xF0, 0xF1)

prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]; W, H = prs.slide_width, prs.slide_height


def tb(s, l, t, w, h, text, size, color=INK, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h)); tf = box.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.text = line; p.alignment = align
        for r in p.runs:
            r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = "Malgun Gothic"
    return box


def title(s, text, tag):
    tb(s, 0.55, 0.32, 11.0, 0.9, text, 28, NAVY, bold=True)
    chip = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.7), Inches(0.38), Inches(1.1), Inches(0.5))
    chip.fill.solid(); chip.fill.fore_color.rgb = {"起": BLUE, "承": GREEN, "轉": ORANGE, "結": PURPLE}.get(tag, GREY)
    chip.line.fill.background()
    tf = chip.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = tag; p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Malgun Gothic"
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2), Inches(12.1), Pt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()


def box(s, l, t, w, h, text, fill, fg=WHITE, size=13):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.color.rgb = WHITE; sh.line.width = Pt(1)
    tf = sh.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.text = line; p.alignment = PP_ALIGN.CENTER
        for r in p.runs:
            r.font.size = Pt(size - (2 if i else 0)); r.font.bold = (i == 0); r.font.color.rgb = fg; r.font.name = "Malgun Gothic"


def arrow(s, l, t, w, color=GREY):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(0.42))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background()


def pic(s, name, l, t, w, cap=""):
    p = IMG / name
    if p.exists():
        ph = s.shapes.add_picture(str(p), Inches(l), Inches(t), width=Inches(w))
        if cap:
            tb(s, l, t + Emu(ph.height).inches + 0.04, w, 0.4, cap, 10.5, GREY, align=PP_ALIGN.CENTER)


# 1 표지
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tb(s, 1.0, 2.3, 11.3, 1.4, "XR 실사급 공간 자산을 위한 SLAM × Gaussian Splatting", 33, WHITE, bold=True)
tb(s, 1.0, 3.8, 11.3, 1.0, "decoupled 파이프라인 — 설계부터 풀맵 실시간 위치추정·실행 통합까지", 19, RGBColor(0xBB, 0xCC, 0xEE))
tb(s, 1.0, 5.0, 11.3, 0.6, "起 목표·한계   承 원리·품질   轉 합쳐진 맵 성능   結 시스템화", 15, RGBColor(0x88, 0xAA, 0xCC))

# 2 起 목표
s = prs.slides.add_slide(BLANK); title(s, "목표", "起")
tb(s, 0.7, 1.6, 12.0, 1.2, "“XR에서 사람이 들어가도 어색하지 않은 실사급 공간 자산”", 24, NAVY, bold=True)
box(s, 1.5, 3.2, 4.2, 1.3, "보이는 화면\n= 실사 같은 3D 공간", BLUE, size=16)
tb(s, 6.0, 3.4, 1.0, 1.0, "+", 40, GREY, bold=True, align=PP_ALIGN.CENTER)
box(s, 7.0, 3.2, 4.2, 1.3, "동시에\n= 내 위치를 안정 추정", GREEN, size=16)
tb(s, 0.7, 5.0, 12.0, 0.8, "핵심 질문: 한 모델로? 아니면 역할을 나눠서? → 나눠서(decoupled)가 답.", 17, ORANGE, bold=True)

# 3 起 통합형 한계
s = prs.slides.add_slide(BLANK); title(s, "통합형 Gaussian-SLAM의 한계", "起")
tb(s, 0.7, 1.6, 12.0, 0.7, "SplaTAM·GS-SLAM·Photo-SLAM·LoopSplat 4종 직접 빌드·실행 → 세 가지가 동시 미달", 17, INK)
box(s, 1.2, 2.7, 3.5, 1.1, "렌더 품질\n실사급 미달", RED, size=15)
box(s, 5.0, 2.7, 3.5, 1.1, "트래킹\n전용 SLAM보다 약함", RED, size=15)
box(s, 8.8, 2.7, 3.5, 1.1, "속도\n실시간 XR 부담", RED, size=15)
tb(s, 0.7, 4.3, 12.0, 1.0, "→ ‘장점 통합’이 아니라 ‘단점 통합’ 위험. 한 표현에 모든 역할을 욱여넣은 대가.", 18, ORANGE, bold=True)

# 4 承 decoupled 전환
s = prs.slides.add_slide(BLANK); title(s, "전환 — 역할을 쪼갠다 (decoupled)", "承")
box(s, 0.8, 2.0, 2.6, 1.2, "D455\nRGB + Depth", BLUE)
arrow(s, 3.5, 2.4, 0.7, GREY); box(s, 4.3, 2.0, 2.6, 1.2, "ORB-SLAM3\n위치 = 좌표계", GREEN)
arrow(s, 7.0, 2.4, 0.7, GREY); box(s, 7.8, 2.0, 2.6, 1.2, "gsplat 학습\n(포즈 고정)", ORANGE)
arrow(s, 10.5, 2.4, 0.6, GREY); box(s, 11.2, 2.0, 1.9, 1.2, "실사 자산\n.ply", PURPLE)
tb(s, 0.8, 3.7, 12.2, 2.2,
   "핵심 트릭: Gaussian을 ‘SLAM 포즈 위에서’ 학습 → 두 맵이 같은 좌표계 → 정합 0단계.\n\n"
   "Localization = 성숙한 ORB-SLAM3 / Rendering = 오프라인 무제약 gsplat.\n"
   "런타임엔 위치추정 포즈를 변환 없이 렌더러에 바로 투입.", 16, GREEN)

# 5 承 M1 원리
s = prs.slides.add_slide(BLANK); title(s, "M1 — 원리 검증 (공개데이터, mocap GT)", "承")
box(s, 1.0, 2.0, 5.4, 1.2, "ORB 포즈 23.85  vs  COLMAP 23.98 PSNR\n→ 차이 Δ0.13 dB", GREEN, size=16)
box(s, 1.0, 3.5, 5.4, 1.2, "ATE vs mocap GT = 1.9 cm\n(유일한 절대 정확도 증명)", BLUE, size=16)
tb(s, 6.8, 2.0, 6.0, 3.0,
   "TUM fr1/desk, hold-out 16뷰.\n\n"
   "‘좋은 입력이면 SLAM 포즈로 COLMAP급\n품질이 나온다’ — decoupled 원리 성립.\n\n"
   "자체 데이터(home)는 mocap이 없어 이\nATE 절대 숫자는 공개데이터에서만.", 16, INK)

# 6 承 품질 MCMC
s = prs.slides.add_slide(BLANK); title(s, "M2 — 실사 품질: 학습 전략(MCMC)으로 +4 dB", "承")
pic(s, "quality_mcmc_AB.png", 0.6, 1.5, 6.4, "좌:GT  가운데:기존 23.82  우:MCMC 27.96 (또렷)")
tb(s, 7.2, 1.7, 5.8, 4.4,
   "home(자체 D455) 자산:\n\n"
   "Default 23.82 → MCMC-2M 27.96 dB\nSSIM 0.785→0.871  LPIPS 0.391→0.267\n\n"
   "발견: 흐릿함은 하드웨어 한계가 아니라\n‘가우시안 배치’ 문제. 같은 2M 예산을\nMCMC로 최적 배치 → 선명.\n\n"
   "통제비교(2M≈3M) → 타일링 불필요.", 15, INK)

# 7 轉 포즈 허용곡선
s = prs.slides.add_slide(BLANK); title(s, "합쳐진 맵 성능 ① — 포즈 허용오차", "轉")
pic(s, "pose_sensitivity_curve.png", 0.6, 1.6, 8.0, "포즈 오차 vs 렌더 화질 (좌:이동, 우:회전)")
tb(s, 8.9, 1.8, 4.1, 4.4,
   "위치추정이 틀린 만큼\n렌더가 나빠진다 = 결합식.\n\n"
   "회전 1° → −7 dB (민감)\n이동 1cm → −2.8 dB\n\n"
   "→ localizer 예산:\n   회전 <1°, 이동 ~1cm.\n\n"
   "(thesis 핵심 숫자)", 15, INK)

# 8 轉 자가 localize
s = prs.slides.add_slide(BLANK); title(s, "합쳐진 맵 성능 ② — 가우시안 맵이 스스로 위치추정", "轉")
pic(s, "photometric_reloc_AB.png", 0.6, 1.5, 6.6, "좌:실제 우:가우시안 렌더(틀어진 포즈→정합)")
tb(s, 7.4, 1.7, 5.6, 4.6,
   "Gaussian 맵을 미분가능 렌더러로:\n5cm/3° 틀어진 포즈 → 이미지에 맞춰\n포즈 최적화 → 0.18cm/0.05°, 10/10.\n\n"
   "수렴 basin: 이동 20cm·회전 12°까지\n75%+ → 거친 localizer 허용.\n\n"
   "‘한 맵이 렌더 + 위치추정 둘 다’.", 15, INK)

# 9 轉 envelope
s = prs.slides.add_slide(BLANK); title(s, "합쳐진 맵 성능 ③ — 동작 범위(envelope)", "轉")
pic(s, "merged_map_summary.png", 0.6, 1.5, 8.4, "5개 결과 종합")
tb(s, 9.2, 1.8, 3.8, 4.4,
   "찍은 공간 안:\nlocalize ✓ · render ✓(28) · photometric ✓\n\n"
   "안 찍은 공간:\nlocalize ✓(0.73cm) · render ✗(10)\n\n"
   "→ 율속 = 캡처 커버리지,\n   위치추정 아님.", 14, INK)

# 10 轉 feature-PnP 실시간
s = prs.slides.add_slide(BLANK); title(s, "합쳐진 맵 성능 ④ — 실시간 위치추정 (feature-PnP)", "轉")
pic(s, "localize_to_render.png", 0.6, 1.5, 5.6, "좌:실제(맵에 없는 새 프레임) 우:찾은 포즈로 렌더")
tb(s, 6.4, 1.7, 6.6, 4.8,
   "새 프레임만 보고 위치를 찾고(전역, hint 없음)\n그 포즈로 풀자산 렌더:\n\n"
   "• 전역 reloc 성공 100% (40/40)\n"
   "• 26.7 FPS (실시간 근처, 광도법 대비 8×)\n"
   "• render-vs-real 28.3 dB\n\n"
   "= ‘위치추정 → 렌더’ 전 과정이 한 좌표계에서\n   닫힌 end-to-end 증거.", 15, INK)

# 11 轉 풀맵 28.8m
s = prs.slides.add_slide(BLANK); title(s, "합쳐진 맵 성능 ⑤ — 풀 28.8m 공간 전역", "轉")
pic(s, "localization_full_map.png", 0.6, 1.5, 7.6, "초록=SLAM 경로(28.8m) 빨강=PnP가 찾은 query (40/40)")
tb(s, 8.5, 1.8, 4.5, 4.4,
   "작은 2m 구간이 아니라\n전체 29m 녹화를 통째로 학습\n→ 풀 28.8m 가우시안 맵.\n\n"
   "non-KF query 40개 전역 localize\n40/40, 28 FPS.\n\n"
   "‘큰 공간 전역에서 작동’ 증명.", 15, INK)

# 12 結 시스템화
s = prs.slides.add_slide(BLANK); title(s, "시스템화 — 자동 게이트 · 런타임 루프 · 단일 CLI", "結")
box(s, 0.8, 1.9, 3.7, 1.5, "자동 자산 게이트\nbuild_report\nXR_READY/FRAME_INVALID\n(PSNR 높아도 프레임 깨지면 거부)", GREEN, size=13)
box(s, 4.8, 1.9, 3.7, 1.5, "런타임 루프\npose stream→localize→render\n상태기계(OK/LOST)\n40/40 트래킹", ORANGE, size=13)
box(s, 8.8, 1.9, 3.7, 1.5, "단일 CLI\nxrsplat build/report\n/view/localize/run\n8단계→1명령(게이트·resume)", PURPLE, size=13)
tb(s, 0.8, 3.8, 12.0, 2.0,
   "파편화된 8스크립트 + 런타임 호출 → config 1개 + 오케스트레이터 + MergedMap API + 단일 현관.\n"
   "‘compose, not rewrite’ — 작동하는 스크립트 유지, 위에 얇은 한 겹. resume이라 학습만 재실행 가능.\n\n"
   "pytest 28 passed · gsplat-free 불변식 유지 · home build = 전단계 skip→XR_READY(idempotent).", 15, INK)

# 13 結 종합 성능표
s = prs.slides.add_slide(BLANK); title(s, "종합 성능 (실측)", "結")
rows = [
    ("원리(M1, GT)", "ATE 1.9cm · ORB vs COLMAP Δ0.13dB", GREEN),
    ("실사 품질(home)", "27.96 dB · SSIM 0.871 · LPIPS 0.267", GREEN),
    ("포즈 허용오차", "회전 1°=−7dB · 이동 1cm=−2.8dB", ORANGE),
    ("자가 위치추정", "5cm/3° → 0.18cm/0.05° · 10/10", ORANGE),
    ("실시간 reloc(PnP)", "전역 100% · 26.7 FPS · 28.3 dB", BLUE),
    ("풀 28.8m 전역", "40/40 localize · 28 FPS", BLUE),
    ("자산 게이트", "home XR_READY · 프레임 0.999/0.05°", PURPLE),
    ("실행 통합", "단일 CLI · pytest 28 · idempotent build", PURPLE),
]
y = 1.55
for name, val, c in rows:
    box(s, 0.8, y, 3.6, 0.62, name, c, size=13)
    tb(s, 4.6, y + 0.05, 8.4, 0.55, val, 15, INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y += 0.7

# 14 結 결과물 갤러리
s = prs.slides.add_slide(BLANK); title(s, "결과물 — 3D 자산(.ply) + 렌더", "結")
pic(s, "home_gallery.png", 4.4, 1.4, 4.4, "home 렌더 7뷰 (27~30 dB)")
tb(s, 0.7, 1.7, 3.4, 4.6,
   "최종 산출물 = 3D Gaussian\n`.ply` (200만 가우시안)\n+ 경량 배포본 scene_lite.ply.\n\n"
   "어느 각도서든 실사 렌더\n(SuperSplat·Unity·Unreal·웹뷰어).\n\n"
   "home 2m + 풀 28.8m 두 자산.", 15, INK)
tb(s, 9.0, 1.7, 4.0, 4.6,
   "왼쪽 7뷰 = 같은 .ply를\n여러 시점에서 렌더한 스냅샷.\n\n"
   "소파·식물·창틀·쿠션 모두\n실제와 충실히 일치.", 15, INK)

# 15 結 한계·다음
s = prs.slides.add_slide(BLANK); title(s, "남은 한계 & 다음 (정직)", "結")
tb(s, 0.8, 1.7, 12.0, 4.4,
   "• 자체 데이터 절대-cm GT 없음(mocap) → render-vs-real + 궤적으로 증명. 절대 ATE는 TUM(M1=1.9cm)만.\n\n"
   "• 실시간 ~27 FPS(30 근접), 오프라인 리플레이 — 실물 HMD/VIO 미연결.\n\n"
   "• 품질 ↔ 범위 trade-off: 풀 28.8m은 밀도 퍼져 soft(22dB), 좁은 2m은 28dB.\n   둘 다 높이려면 더 촘촘히 캡처/학습 (알고리즘 아닌 데이터·연산 문제).\n\n"
   "다음 단계 = ‘제품화’: 실물 HMD 연결 · 절대 GT 검증 · 커버리지 확대.", 16, INK)

# 16 結 결론
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H); bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
tb(s, 1.0, 2.2, 11.3, 2.0,
   "통합 모델 버리고 SLAM(위치)+Gaussian(렌더)을 분리,\n같은 좌표계로 연동 → XR 실사 공간이 실제로 작동한다.", 26, WHITE, bold=True)
tb(s, 1.0, 4.6, 11.3, 1.4,
   "설계 → 자산 → 합쳐진 맵 성능 → 풀맵 실시간 위치추정 → 실행 통합까지 전 구간 실증 완료.",
   18, RGBColor(0xBB, 0xCC, 0xEE))

OUT.parent.mkdir(parents=True, exist_ok=True); prs.save(str(OUT))
print(f"[final] {len(prs.slides._sldIdLst)} slides → {OUT}")
